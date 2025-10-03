import imaplib
import email
from email.header import decode_header
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import re
from bs4 import BeautifulSoup
import quopri
import base64
import os
from pathlib import Path
from datetime import datetime
import mimetypes

# 첨부파일 텍스트 추출용 라이브러리
try:
    import PyPDF2  # PDF 파일
except ImportError:
    PyPDF2 = None

try:
    from docx import Document  # Word 파일
except ImportError:
    Document = None

try:
    import openpyxl  # Excel 파일
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation  # PowerPoint 파일
except ImportError:
    Presentation = None

class NaverMailParser:
    def __init__(self, username, password, download_path='downloads'):
        self.username = username
        self.password = password
        self.mail = None
        self.download_path = download_path

        # 다운로드 디렉토리 생성
        Path(self.download_path).mkdir(exist_ok=True)

    def connect(self):
        """네이버 IMAP 서버에 연결"""
        try:
            self.mail = imaplib.IMAP4_SSL('imap.naver.com', 993)
            self.mail.login(self.username, self.password)
            return True
        except Exception as e:
            print(f"연결 실패: {e}")
            return False

    def select_mailbox(self, mailbox='INBOX'):
        """메일박스 선택"""
        if self.mail:
            self.mail.select(mailbox)

    def search_emails(self, criteria='ALL', limit=10):
        """이메일 검색"""
        if not self.mail:
            print("먼저 연결해주세요.")
            return []

        try:
            status, messages = self.mail.search(None, criteria)
            if status == 'OK':
                email_ids = messages[0].split()
                return email_ids[-limit:][::-1]  # 최신 메일부터 정렬
        except Exception as e:
            print(f"검색 실패: {e}")
        return []

    def decode_mime_words(self, text):
        """MIME 인코딩된 텍스트 디코딩"""
        if text is None:
            return ""

        decoded_fragments = decode_header(text)
        decoded_text = ""

        for fragment, encoding in decoded_fragments:
            if isinstance(fragment, bytes):
                if encoding:
                    try:
                        decoded_text += fragment.decode(encoding)
                    except:
                        decoded_text += fragment.decode('utf-8', errors='ignore')
                else:
                    decoded_text += fragment.decode('utf-8', errors='ignore')
            else:
                decoded_text += str(fragment)

        return decoded_text

    def extract_text_from_html(self, html_content):
        """HTML에서 텍스트 추출"""
        try:
            soup = BeautifulSoup(html_content, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
        except:
            return html_content

    def sanitize_filename(self, filename):
        """파일명에서 특수문자 제거"""
        # 위험한 문자 제거
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # 공백을 언더스코어로
        filename = filename.replace(' ', '_')
        # 길이 제한
        if len(filename) > 200:
            name, ext = os.path.splitext(filename)
            filename = name[:200-len(ext)] + ext
        return filename

    def extract_text_from_file(self, filepath):
        """첨부파일에서 텍스트 추출"""
        try:
            # 파일 확장자 확인
            _, ext = os.path.splitext(filepath)
            ext = ext.lower()

            # PDF 파일
            if ext == '.pdf' and PyPDF2:
                try:
                    with open(filepath, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        text_parts = []
                        for page_num in range(len(pdf_reader.pages)):
                            page = pdf_reader.pages[page_num]
                            text_parts.append(page.extract_text())
                        return '\n'.join(text_parts)
                except Exception as e:
                    return f"[PDF 텍스트 추출 실패: {e}]"

            # Word 문서
            elif ext in ['.docx', '.doc'] and Document:
                try:
                    doc = Document(filepath)
                    text_parts = []
                    for paragraph in doc.paragraphs:
                        text_parts.append(paragraph.text)
                    return '\n'.join(text_parts)
                except Exception as e:
                    return f"[Word 텍스트 추출 실패: {e}]"

            # Excel 파일
            elif ext in ['.xlsx', '.xls'] and openpyxl:
                try:
                    wb = openpyxl.load_workbook(filepath, data_only=True)
                    text_parts = []
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        text_parts.append(f"[Sheet: {sheet_name}]")
                        for row in sheet.iter_rows(values_only=True):
                            row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                            if row_text.strip():
                                text_parts.append(row_text)
                    return '\n'.join(text_parts)
                except Exception as e:
                    return f"[Excel 텍스트 추출 실패: {e}]"

            # PowerPoint 파일
            elif ext in ['.pptx', '.ppt'] and Presentation:
                try:
                    prs = Presentation(filepath)
                    text_parts = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text_parts.append(f"[Slide {slide_num}]")
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_parts.append(shape.text)
                    return '\n'.join(text_parts)
                except Exception as e:
                    return f"[PowerPoint 텍스트 추출 실패: {e}]"

            # 텍스트 파일
            elif ext in ['.txt', '.csv', '.log', '.md', '.json', '.xml', '.html', '.htm']:
                try:
                    # 다양한 인코딩 시도
                    encodings = ['utf-8', 'euc-kr', 'cp949', 'latin-1']
                    for encoding in encodings:
                        try:
                            with open(filepath, 'r', encoding=encoding) as f:
                                return f.read()
                        except UnicodeDecodeError:
                            continue
                    return "[텍스트 파일 인코딩 실패]"
                except Exception as e:
                    return f"[텍스트 파일 읽기 실패: {e}]"

            # 지원하지 않는 형식
            else:
                return f"[{ext} 형식은 텍스트 추출을 지원하지 않습니다]"

        except Exception as e:
            return f"[파일 처리 오류: {e}]"

    def save_attachment(self, part, email_id, extract_text=False):
        """첨부파일 저장 및 텍스트 추출"""
        try:
            filename = part.get_filename()
            if filename:
                # 파일명 디코딩
                filename = self.decode_mime_words(filename)
                filename = self.sanitize_filename(filename)

                # 이메일 ID별 폴더 생성
                email_folder = Path(self.download_path) / f"email_{email_id}"
                email_folder.mkdir(exist_ok=True)

                filepath = email_folder / filename

                # 중복 파일명 처리
                counter = 1
                original_filepath = filepath
                while filepath.exists():
                    name, ext = os.path.splitext(original_filepath)
                    filepath = Path(f"{name}_{counter}{ext}")
                    counter += 1

                # 첨부파일 저장
                with open(filepath, 'wb') as f:
                    f.write(part.get_payload(decode=True))

                attachment_info = {
                    'filename': filename,
                    'filepath': str(filepath),
                    'size': os.path.getsize(filepath),
                    'content_type': part.get_content_type()
                }

                # 텍스트 추출 옵션이 활성화된 경우
                if extract_text:
                    extracted_text = self.extract_text_from_file(str(filepath))
                    attachment_info['extracted_text'] = extracted_text

                return attachment_info
        except Exception as e:
            print(f"첨부파일 저장 실패: {e}")
        return None

    def parse_email_body_and_attachments(self, msg, email_id, extract_attachment_text=True):
        """이메일 본문과 첨부파일 파싱 (첨부파일 텍스트 추출 포함)"""
        body_text = ""
        body_html = ""
        attachments = []
        attachment_texts = []

        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                content_disposition = str(part.get("Content-Disposition"))

                # 첨부파일 처리
                if "attachment" in content_disposition:
                    attachment_info = self.save_attachment(part, email_id, extract_text=extract_attachment_text)
                    if attachment_info:
                        attachments.append(attachment_info)
                        # 추출된 텍스트가 있으면 저장
                        if extract_attachment_text and 'extracted_text' in attachment_info:
                            attachment_texts.append(f"\n\n=== 첨부파일: {attachment_info['filename']} ===\n{attachment_info['extracted_text']}")

                # 본문 처리
                else:
                    if content_type == "text/plain":
                        charset = part.get_content_charset() or 'utf-8'
                        payload = part.get_payload(decode=True)
                        if payload:
                            try:
                                body_text += payload.decode(charset, errors='ignore')
                            except:
                                body_text += payload.decode('utf-8', errors='ignore')

                    elif content_type == "text/html":
                        charset = part.get_content_charset() or 'utf-8'
                        payload = part.get_payload(decode=True)
                        if payload:
                            try:
                                body_html += payload.decode(charset, errors='ignore')
                            except:
                                body_html += payload.decode('utf-8', errors='ignore')
        else:
            content_type = msg.get_content_type()
            charset = msg.get_content_charset() or 'utf-8'
            payload = msg.get_payload(decode=True)

            if payload:
                try:
                    content = payload.decode(charset, errors='ignore')
                except:
                    content = payload.decode('utf-8', errors='ignore')

                if content_type == "text/html":
                    body_html = content
                else:
                    body_text = content

        # HTML이 있으면 텍스트로 변환
        if body_html and not body_text:
            body_text = self.extract_text_from_html(body_html)

        # 첨부파일 텍스트를 본문에 추가
        if attachment_texts:
            combined_text = body_text.strip() + '\n'.join(attachment_texts)
        else:
            combined_text = body_text.strip()

        return {
            'body_text': combined_text,
            'body_html': body_html.strip(),
            'attachments': attachments,
            'has_extracted_attachments': len(attachment_texts) > 0
        }

    def download_email_full(self, email_id, extract_attachment_text=True):
        """전체 이메일 다운로드 (본문 + 첨부파일 + 첨부파일 텍스트 추출)"""
        if not self.mail:
            return None

        try:
            status, msg_data = self.mail.fetch(email_id, '(RFC822)')
            if status != 'OK':
                return None

            raw_email = msg_data[0][1]
            msg = email.message_from_bytes(raw_email)

            # 헤더 정보 추출
            subject = self.decode_mime_words(msg.get("Subject", ""))
            sender = self.decode_mime_words(msg.get("From", ""))
            recipient = self.decode_mime_words(msg.get("To", ""))
            date = msg.get("Date", "")

            # 본문과 첨부파일 추출 (텍스트 추출 포함)
            email_id_str = email_id.decode() if isinstance(email_id, bytes) else str(email_id)
            parsed_content = self.parse_email_body_and_attachments(msg, email_id_str, extract_attachment_text=extract_attachment_text)

            # 이메일 전체 정보 저장
            email_folder = Path(self.download_path) / f"email_{email_id_str}"
            email_folder.mkdir(exist_ok=True)

            # 메타데이터 저장
            metadata = {
                'id': email_id_str,
                'subject': subject,
                'sender': sender,
                'recipient': recipient,
                'date': date,
                'downloaded_at': datetime.now().isoformat(),
                'has_html': bool(parsed_content['body_html']),
                'attachment_count': len(parsed_content['attachments'])
            }

            # 본문 텍스트 저장
            if parsed_content['body_text']:
                text_file = email_folder / 'body.txt'
                with open(text_file, 'w', encoding='utf-8') as f:
                    f.write(parsed_content['body_text'])

            # HTML 본문 저장
            if parsed_content['body_html']:
                html_file = email_folder / 'body.html'
                with open(html_file, 'w', encoding='utf-8') as f:
                    f.write(parsed_content['body_html'])

            # 메타데이터 저장
            import json
            metadata_file = email_folder / 'metadata.json'
            with open(metadata_file, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, ensure_ascii=False, indent=2)

            return {
                'id': email_id_str,
                'subject': subject,
                'sender': sender,
                'recipient': recipient,
                'date': date,
                'body': parsed_content['body_text'],
                'body_html': parsed_content['body_html'],
                'attachments': parsed_content['attachments'],
                'download_folder': str(email_folder)
            }

        except Exception as e:
            print(f"이메일 다운로드 실패 (ID: {email_id}): {e}")
            return None

    def parse_email(self, email_id, extract_attachment_text=False):
        """개별 이메일 파싱 (기본 버전)"""
        if not self.mail:
            return None

        try:
            status, msg_data = self.mail.fetch(email_id, '(RFC822)')
            if status != 'OK':
                return None

            raw_email = msg_data[0][1]
            msg = email.message_from_bytes(raw_email)

            # 헤더 정보 추출
            subject = self.decode_mime_words(msg.get("Subject", ""))
            sender = self.decode_mime_words(msg.get("From", ""))
            recipient = self.decode_mime_words(msg.get("To", ""))
            date = msg.get("Date", "")

            # 본문 추출
            email_id_str = email_id.decode() if isinstance(email_id, bytes) else str(email_id)
            parsed_content = self.parse_email_body_and_attachments(msg, email_id_str, extract_attachment_text=extract_attachment_text)

            # 첨부파일 개수 확인
            attachment_count = len(parsed_content['attachments'])

            return {
                'id': email_id_str,
                'subject': subject,
                'sender': sender,
                'recipient': recipient,
                'date': date,
                'body': parsed_content['body_text'],
                'has_attachments': attachment_count > 0,
                'attachment_count': attachment_count
            }

        except Exception as e:
            print(f"이메일 파싱 실패 (ID: {email_id}): {e}")
            return None

    def get_emails(self, criteria='ALL', limit=10, download_full=False, extract_attachment_text=True):
        """이메일 목록 가져오기"""
        email_ids = self.search_emails(criteria, limit)
        emails = []

        for email_id in email_ids:
            if download_full:
                parsed_email = self.download_email_full(email_id, extract_attachment_text=extract_attachment_text)
            else:
                parsed_email = self.parse_email(email_id, extract_attachment_text=extract_attachment_text)

            if parsed_email:
                emails.append(parsed_email)

        return emails

    def view_email_content(self, email_id):
        """특정 이메일의 전체 내용 보기"""
        email_folder = Path(self.download_path) / f"email_{email_id}"

        if not email_folder.exists():
            # 다운로드되지 않은 경우 새로 다운로드
            email_id_bytes = email_id.encode() if isinstance(email_id, str) else email_id
            return self.download_email_full(email_id_bytes)

        # 이미 다운로드된 경우 파일에서 읽기
        try:
            import json

            # 메타데이터 읽기
            metadata_file = email_folder / 'metadata.json'
            with open(metadata_file, 'r', encoding='utf-8') as f:
                metadata = json.load(f)

            # 본문 읽기
            body_text = ""
            text_file = email_folder / 'body.txt'
            if text_file.exists():
                with open(text_file, 'r', encoding='utf-8') as f:
                    body_text = f.read()

            body_html = ""
            html_file = email_folder / 'body.html'
            if html_file.exists():
                with open(html_file, 'r', encoding='utf-8') as f:
                    body_html = f.read()

            # 첨부파일 목록
            attachments = []
            for file in email_folder.iterdir():
                if file.name not in ['body.txt', 'body.html', 'metadata.json']:
                    attachments.append({
                        'filename': file.name,
                        'filepath': str(file),
                        'size': file.stat().st_size
                    })

            return {
                **metadata,
                'body': body_text,
                'body_html': body_html,
                'attachments': attachments,
                'download_folder': str(email_folder)
            }

        except Exception as e:
            print(f"이메일 내용 읽기 실패: {e}")
            return None

    def get_attachment_path(self, email_id, filename):
        """특정 첨부파일의 경로 반환"""
        email_folder = Path(self.download_path) / f"email_{email_id}"
        attachment_path = email_folder / filename

        if attachment_path.exists():
            return str(attachment_path)
        return None

    def close(self):
        """연결 종료"""
        if self.mail:
            self.mail.close()
            self.mail.logout()

def main():
    # 사용 예시
    username = "your_naver_id@naver.com"
    password = "your_app_password"  # 네이버 앱 비밀번호 사용 권장

    parser = NaverMailParser(username, password)

    if parser.connect():
        print("네이버 메일에 연결되었습니다.")

        # 받은편지함 선택
        parser.select_mailbox('INBOX')

        # 최근 5개 이메일 가져오기 (첨부파일 포함)
        print("\n=== 이메일 다운로드 중 ===")
        emails = parser.get_emails(limit=5, download_full=True)

        for i, email_data in enumerate(emails, 1):
            print(f"\n=== 이메일 {i} ===")
            print(f"제목: {email_data['subject']}")
            print(f"발신자: {email_data['sender']}")
            print(f"날짜: {email_data['date']}")
            print(f"첨부파일: {email_data.get('attachment_count', 0)}개")

            if email_data.get('attachments'):
                print("첨부파일 목록:")
                for att in email_data['attachments']:
                    print(f"  - {att['filename']} ({att['size']} bytes)")

            print(f"본문 (처음 200자):")
            body = email_data.get('body', '')
            print(body[:200] + "..." if len(body) > 200 else body)
            print(f"다운로드 폴더: {email_data.get('download_folder', 'N/A')}")

        parser.close()
    else:
        print("네이버 메일 연결에 실패했습니다.")

if __name__ == "__main__":
    main()